"""
Office文档解析器
支持Word、PowerPoint、Excel文档的解析
"""
import asyncio
import time
from pathlib import Path
from typing import List, Dict, Any, Optional
import json

from docx import Document as DocxDocument
from docx.shared import Inches
from pptx import Presentation
import pandas as pd
import openpyxl

from src.data.parsers.base import BaseParser, ParseResult, ParsedElement, MultiModalParser
from src.data.document_validator import DocumentInfo, DocumentType
from src.utils.logger import get_logger
from config.settings import settings

logger = get_logger("data.parsers.office")


class WordParser(BaseParser):
    """Word文档解析器"""
    
    def __init__(self):
        super().__init__("WordParser")
        self.supported_types = [DocumentType.DOCX, DocumentType.DOC]
    
    async def parse(self, file_path: Path, doc_info: DocumentInfo) -> ParseResult:
        """解析Word文档"""
        start_time = time.time()
        
        try:
            doc = DocxDocument(str(file_path))
            elements = []
            
            # 解析段落
            for para_index, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_element = self._create_text_element(
                        content=paragraph.text,
                        metadata={
                            'paragraph_index': para_index,
                            'style': paragraph.style.name if paragraph.style else None,
                            'extraction_method': 'python_docx'
                        },
                        position={'paragraph': para_index}
                    )
                    elements.append(text_element)
            
            # 解析表格
            for table_index, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                markdown_table = self._convert_table_to_markdown(table_data)
                table_element = self._create_table_element(
                    content=markdown_table,
                    metadata={
                        'table_index': table_index,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0,
                        'extraction_method': 'python_docx'
                    },
                    position={'table': table_index}
                )
                elements.append(table_element)
            
            # 解析图片（获取图片信息）
            image_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_element = self._create_image_element(
                        content=f"[图像 {image_count + 1}]",
                        metadata={
                            'image_index': image_count,
                            'image_path': rel.target_ref,
                            'extraction_method': 'python_docx'
                        },
                        position={'image': image_count}
                    )
                    elements.append(image_element)
                    image_count += 1
            
            processing_time = time.time() - start_time
            
            return ParseResult(
                success=True,
                elements=elements,
                metadata={
                    'total_paragraphs': len(doc.paragraphs),
                    'total_tables': len(doc.tables),
                    'total_images': image_count,
                    'parser': self.name,
                    'extraction_method': 'python_docx'
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            return self._handle_error(e, file_path)
    
    def _convert_table_to_markdown(self, table_data: List[List[str]]) -> str:
        """将表格数据转换为Markdown格式"""
        if not table_data:
            return ""
        
        markdown_lines = []
        
        # 表头
        if table_data:
            header = "| " + " | ".join(str(cell or "") for cell in table_data[0]) + " |"
            markdown_lines.append(header)
            
            # 分隔线
            separator = "| " + " | ".join("---" for _ in table_data[0]) + " |"
            markdown_lines.append(separator)
            
            # 数据行
            for row in table_data[1:]:
                row_line = "| " + " | ".join(str(cell or "") for cell in row) + " |"
                markdown_lines.append(row_line)
        
        return "\n".join(markdown_lines)


class PowerPointParser(BaseParser):
    """PowerPoint文档解析器"""
    
    def __init__(self):
        super().__init__("PowerPointParser")
        self.supported_types = [DocumentType.PPTX, DocumentType.PPT]
    
    async def parse(self, file_path: Path, doc_info: DocumentInfo) -> ParseResult:
        """解析PowerPoint文档"""
        start_time = time.time()
        
        try:
            prs = Presentation(str(file_path))
            elements = []
            
            for slide_index, slide in enumerate(prs.slides):
                # 解析文本
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if slide_texts:
                    combined_text = "\n".join(slide_texts)
                    text_element = self._create_text_element(
                        content=combined_text,
                        metadata={
                            'slide_index': slide_index,
                            'slide_number': slide_index + 1,
                            'extraction_method': 'python_pptx'
                        },
                        position={'slide': slide_index + 1}
                    )
                    elements.append(text_element)
                
                # 解析表格
                table_count = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        markdown_table = self._convert_table_to_markdown(table_data)
                        table_element = self._create_table_element(
                            content=markdown_table,
                            metadata={
                                'slide_index': slide_index,
                                'slide_number': slide_index + 1,
                                'table_index': table_count,
                                'rows': len(table_data),
                                'cols': len(table_data[0]) if table_data else 0,
                                'extraction_method': 'python_pptx'
                            },
                            position={'slide': slide_index + 1, 'table': table_count}
                        )
                        elements.append(table_element)
                        table_count += 1
                
                # 解析图片
                image_count = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture type
                        image_element = self._create_image_element(
                            content=f"[图像 {image_count + 1} 在幻灯片 {slide_index + 1}]",
                            metadata={
                                'slide_index': slide_index,
                                'slide_number': slide_index + 1,
                                'image_index': image_count,
                                'extraction_method': 'python_pptx'
                            },
                            position={'slide': slide_index + 1, 'image': image_count}
                        )
                        elements.append(image_element)
                        image_count += 1
            
            processing_time = time.time() - start_time
            
            return ParseResult(
                success=True,
                elements=elements,
                metadata={
                    'total_slides': len(prs.slides),
                    'parser': self.name,
                    'extraction_method': 'python_pptx'
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            return self._handle_error(e, file_path)
    
    def _convert_table_to_markdown(self, table_data: List[List[str]]) -> str:
        """将表格数据转换为Markdown格式"""
        if not table_data:
            return ""
        
        markdown_lines = []
        
        # 表头
        if table_data:
            header = "| " + " | ".join(str(cell or "") for cell in table_data[0]) + " |"
            markdown_lines.append(header)
            
            # 分隔线
            separator = "| " + " | ".join("---" for _ in table_data[0]) + " |"
            markdown_lines.append(separator)
            
            # 数据行
            for row in table_data[1:]:
                row_line = "| " + " | ".join(str(cell or "") for cell in row) + " |"
                markdown_lines.append(row_line)
        
        return "\n".join(markdown_lines)


class ExcelParser(BaseParser):
    """Excel文档解析器"""
    
    def __init__(self):
        super().__init__("ExcelParser")
        self.supported_types = [DocumentType.XLSX, DocumentType.XLS]
    
    async def parse(self, file_path: Path, doc_info: DocumentInfo) -> ParseResult:
        """解析Excel文档"""
        start_time = time.time()
        
        try:
            # 使用pandas读取Excel文件
            excel_file = pd.ExcelFile(str(file_path))
            elements = []
            
            for sheet_index, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 跳过空工作表
                if df.empty:
                    continue
                
                # 转换为Markdown表格
                markdown_table = df.to_markdown(index=False)
                
                table_element = self._create_table_element(
                    content=markdown_table,
                    metadata={
                        'sheet_index': sheet_index,
                        'sheet_name': sheet_name,
                        'rows': len(df),
                        'cols': len(df.columns),
                        'extraction_method': 'pandas'
                    },
                    position={'sheet': sheet_index, 'sheet_name': sheet_name}
                )
                elements.append(table_element)
                
                # 检查公式（使用openpyxl）
                try:
                    wb = openpyxl.load_workbook(str(file_path))
                    ws = wb[sheet_name]
                    
                    formulas = []
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.data_type == 'f':  # 公式类型
                                formulas.append({
                                    'cell': cell.coordinate,
                                    'formula': cell.value
                                })
                    
                    if formulas:
                        formula_text = "\n".join([f"{f['cell']}: {f['formula']}" for f in formulas])
                        formula_element = self._create_formula_element(
                            content=formula_text,
                            metadata={
                                'sheet_index': sheet_index,
                                'sheet_name': sheet_name,
                                'formula_count': len(formulas),
                                'extraction_method': 'openpyxl'
                            },
                            position={'sheet': sheet_index, 'sheet_name': sheet_name}
                        )
                        elements.append(formula_element)
                    
                except Exception as e:
                    logger.warning(f"公式提取失败 (工作表 {sheet_name}): {e}")
            
            processing_time = time.time() - start_time
            
            return ParseResult(
                success=True,
                elements=elements,
                metadata={
                    'total_sheets': len(excel_file.sheet_names),
                    'sheet_names': excel_file.sheet_names,
                    'parser': self.name,
                    'extraction_method': 'pandas_openpyxl'
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            return self._handle_error(e, file_path)
